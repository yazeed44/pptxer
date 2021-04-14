import json
import os

from pptx import Presentation

from presentations_scraper import scrape_presentations_to_dir
from util import calculate_length_stats_for_list_of_strings, load_config

config = load_config()


# TODO write documentation for functions
def load_presentations_objects_from_dir(dir_path):
    presentations = []
    for root, dirs, files in os.walk(dir_path):
        for name in files:
            if name.endswith("pptx"):
                path = os.path.join(root, name)
                try:
                    presentations.append({"path": path, "presentationObj": Presentation(path)})
                except:
                    print(f"Unable to process {path} . It is likely to be corrupted or incomplete")
    return presentations


def extract_presentation_text(presentation, is_flattened):
    # Collect notes
    slides_notes_texts = [
        slide.notes_slide.notes_text_frame.text if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None
        else "" for slide in presentation["presentationObj"].slides]
    # Collect text
    slides_body_texts = ["".join([shape.text if shape.has_text_frame else "" for shape in slide.shapes]) for slide in
                         presentation["presentationObj"].slides]

    result = {"path": presentation["path"]}
    if is_flattened:
        notes_text_flattened = {f"slide{i}NoteText": slides_notes_texts[i] for i in range(len(slides_notes_texts))}
        body_text_flattened = {f"slide{i}BodyText": slides_body_texts[i] for i in range(len(slides_body_texts))}
        return {**result,
                **body_text_flattened,
                **notes_text_flattened,
                **calculate_length_stats_for_list_of_strings(slides_body_texts, "BodyText"),
                **calculate_length_stats_for_list_of_strings(slides_notes_texts, "NotesText"), }
    else:
        return {**result, "slides": [{"noteText": note, "bodyText": body} for note, body in
                                     zip(slides_notes_texts, slides_body_texts)],
                "bodyTextLengthStats": calculate_length_stats_for_list_of_strings(slides_body_texts),
                "noteTextLengthStats": calculate_length_stats_for_list_of_strings(slides_notes_texts)}


def extract_presentations_texts(dir_path=config["downloadDirectory"], flattened=False,
                                result_json_file_path=""):
    presentations = load_presentations_objects_from_dir(dir_path)
    presentations = [extract_presentation_text(presentation, flattened) for presentation in
                     presentations]
    if result_json_file_path != "":
        with open(result_json_file_path, 'w') as f:
            json.dump(presentations, f, ensure_ascii=False)
    return presentations


def main():
    # presentation_urls = scrape_presentation_urls()
    # presentation_urls = ["https://enable.unc.edu/wp-content/uploads/2019/05/CHIP_XRHealthcare.pptx"]

    print(extract_presentations_texts(result_json_file_path="parsing_presentations_result.json"))


if __name__ == '__main__':
    main()
