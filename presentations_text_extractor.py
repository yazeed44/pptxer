import json
import os

from pptx import Presentation

from util import calculate_length_stats_for_list_of_strings, load_config

# TODO write documentation for functions


config = load_config()


def extract_presentations_texts(paths, combine_results_of_different_paths=True,
                                presentation_slide_texts_flattened=False, result_json_file_path=""):
    if isinstance(paths, str):
        presentations = __extract_presentation_texts_from_path__(paths, presentation_slide_texts_flattened)
    else:
        presentations = [__extract_presentation_texts_from_path__(path, presentation_slide_texts_flattened) for path in
                         paths]

        presentations = sum(presentations, []) if combine_results_of_different_paths else presentations

    if result_json_file_path != "":
        with open(result_json_file_path, 'w') as f:
            json.dump(presentations, f, ensure_ascii=False)

    return presentations


def __extract_presentation_texts_from_path__(path, flattened, ):
    # If path is to a directory, then navigate that directory and extract text from any pptx file within it
    # If path is to a single file, then extract text from that specific pptx file
    if os.path.isdir(path):
        presentations = __load_presentations_objects_from_dir__(path)
    else:
        presentations = __load_presentations_objects_from_file_paths__(path)
    presentations = [__extract_presentation_text__(presentation, flattened) for presentation in
                     presentations]

    return presentations


def __load_presentations_objects_from_file_paths__(file_paths):
    if isinstance(file_paths, str):
        file_paths = [file_paths]
    presentations = []
    for path in file_paths:
        try:
            presentations.append({"path": path, "presentationObj": Presentation(path)})
        except:
            print(
                f"Unable to process {path} . It is likely to be corrupted or incomplete. Please ensure that the input is a valid pptx file")
    return presentations


def __load_presentations_objects_from_dir__(dir_path):
    paths = []
    for root, dirs, files in os.walk(dir_path):
        for name in files:
            if name.endswith("pptx"):
                paths.append(os.path.join(root, name))

    return __load_presentations_objects_from_file_paths__(paths)


def __extract_presentation_text__(presentation, is_flattened):
    # Collect notes
    slides_notes_texts = [
        slide.notes_slide.notes_text_frame.text if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None
        else "" for slide in presentation["presentationObj"].slides]
    # Collect text
    slides_body_texts = ["".join([shape.text if shape.has_text_frame else "" for shape in slide.shapes]) for slide in
                         presentation["presentationObj"].slides]

    result = {"path": presentation["path"]}
    if is_flattened:
        notes_text_flattened = {f"slide{i}NoteText": slides_notes_texts[i] for i in range(len(slides_notes_texts))}
        body_text_flattened = {f"slide{i}BodyText": slides_body_texts[i] for i in range(len(slides_body_texts))}
        return {**result,
                **body_text_flattened,
                **notes_text_flattened,
                **calculate_length_stats_for_list_of_strings(slides_body_texts, "BodyText"),
                **calculate_length_stats_for_list_of_strings(slides_notes_texts, "NotesText"), }
    else:
        return {**result, "slides": [{"noteText": note, "bodyText": body} for note, body in
                                     zip(slides_notes_texts, slides_body_texts)],
                "bodyTextLengthStats": calculate_length_stats_for_list_of_strings(slides_body_texts),
                "noteTextLengthStats": calculate_length_stats_for_list_of_strings(slides_notes_texts)}
