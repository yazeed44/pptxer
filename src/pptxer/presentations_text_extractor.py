"""This module extract text from presentation files (pptx)"""
import json
import os
import time
from typing import List

from pptx import Presentation
from pptx.exc import PackageNotFoundError

from pptxer.util import __calculate_length_stats_for_list_of_strings__
from pptxer import logger


def extract_presentations_texts(
        paths: List[str],
        single_array_result=True,
        text_fields_flattened=False,
        extract_output_file_path=None,
) -> List[dict]:
    """
    Takes a list of paths to pptx files and returns a list of dictionary with each presentation slides' texts.

            Parameters:
                    paths (List[str]): Paths to pptx files
                    single_array_result (bool): if true then results of extracting texts from different paths
                    will be combined into one array. Otherwise, Have 2D lists of each path result
                    text_fields_flattened (bool): Flatten all text extractions to be one level
                    extract_output_file_path (str): File path to which output (json) will be stored

            Returns:
                    presentations_texts (List[dict]): A list of dict with following attributes:
                    {
                        'path': '$path_to_file',
                        'slides':
                        [
                            {'noteText': $slide_1_note_text, 'bodyText' : $slide_1_body_text}, ...
                            'bodyTextLengthStats': {'totalLength': $total_length_of_body_texts,
                            'avgLength': $avg_length_of_body_texts, 'minLength': $min_length_of_body_texts,
                            'maxLength': $max_length_of_body_texts, 'medianLength': $median_length_of_body_texts},
                            'noteTextLengthStats': {'totalLength': $total_length_of_note_texts,
                            'avgLength': $avg_length_of_note_texts, 'minLength': $min_length_of_note_texts,
                            'maxLength': $max_length_of_note_texts, 'medianLength': $median_length_of_note_texts}
                        ]
                    }
    """
    if extract_output_file_path is None:
        extract_output_file_path = f"presentations_text_{time.time()}.json"
        logger.info(
            "Since no extract output was specified, the output will be "
            "written to => %s",
            extract_output_file_path,
        )
    logger.info(
        "Will start extracting presentation texts with following parameters:"
        "paths = %s, single_array_result = %s,"
        "text_fields_flattened = %s, "
        "extract_output_file_path = %s",
        paths,
        single_array_result,
        text_fields_flattened,
        extract_output_file_path,
    )
    if isinstance(paths, str):
        presentations = __extract_presentation_texts_from_path__(
            paths, text_fields_flattened
        )
    else:
        presentations = [
            __extract_presentation_texts_from_path__(path, text_fields_flattened)
            for path in paths
        ]

        presentations = sum(presentations, []) if single_array_result else presentations

    with open(extract_output_file_path, "w") as output_file:
        json.dump(presentations, output_file, ensure_ascii=False)

    return presentations


def __extract_presentation_texts_from_path__(path: str, flattened: bool):
    if os.path.isdir(path):
        logger.debug("%s is to a dir. Will attempt to all pptx files within it", path)
        presentations = __load_presentations_objects_from_dir__(path)
    else:
        presentations = __load_presentations_objects_from_file_paths__(path)
    presentations = [
        __extract_presentation_text__(presentation, flattened)
        for presentation in presentations
    ]

    return presentations


def __load_presentations_objects_from_file_paths__(file_paths):
    if isinstance(file_paths, str):
        file_paths = [file_paths]
    presentations = []
    for path in file_paths:
        if not os.path.exists(path):
            logger.error("File %s does not exist", path)
            logger.info("Will skip %s due to an error", path)
            continue
        try:
            obj = {"path": path, "presentationObj": Presentation(path)}

            presentations.append(obj)
            logger.debug("Loaded %s successfully", path)
        except (ModuleNotFoundError, PackageNotFoundError):
            logger.error(
                "Unable to process %s . It is likely to be corrupted or incomplete. "
                "Please ensure the input is a valid pptx file",
                path,
            )
            logger.info("Will skip %s due to an error", path)
    return presentations


def __load_presentations_objects_from_dir__(dir_path: str):
    paths = []
    for root, _, files in os.walk(dir_path):
        for name in files:
            if name.endswith("pptx"):
                paths.append(os.path.join(root, name))

    return __load_presentations_objects_from_file_paths__(paths)


def __extract_presentation_text__(presentation: Presentation, is_flattened: bool):
    # Collect notes
    slides_notes_texts = [
        slide.notes_slide.notes_text_frame.text
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None
        else ""
        for slide in presentation["presentationObj"].slides
    ]
    # Collect text
    slides_body_texts = [
        "".join([shape.text if shape.has_text_frame else "" for shape in slide.shapes])
        for slide in presentation["presentationObj"].slides
    ]

    result = {"path": presentation["path"]}
    if is_flattened:
        notes_text_flattened = {
            f"slide{i}NoteText": slides_notes_texts[i]
            for i in range(len(slides_notes_texts))
        }
        body_text_flattened = {
            f"slide{i}BodyText": slides_body_texts[i]
            for i in range(len(slides_body_texts))
        }
        return {
            **result,
            **body_text_flattened,
            **notes_text_flattened,
            **__calculate_length_stats_for_list_of_strings__(
                slides_body_texts, "BodyText"
            ),
            **__calculate_length_stats_for_list_of_strings__(
                slides_notes_texts, "NotesText"
            ),
        }

    return {
        **result,
        "slides": [
            {"noteText": note, "bodyText": body}
            for note, body in zip(slides_notes_texts, slides_body_texts)
        ],
        "bodyTextLengthStats": __calculate_length_stats_for_list_of_strings__(
            slides_body_texts
        ),
        "noteTextLengthStats": __calculate_length_stats_for_list_of_strings__(
            slides_notes_texts
        ),
    }
