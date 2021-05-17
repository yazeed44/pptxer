import json
import logging
import os
import time
from typing import List

from pptx import Presentation

from util import calculate_length_stats_for_list_of_strings


def extract_presentations_texts(
        paths,
        single_array_result=True,
        text_fields_flattened=False,
        extract_output_file_path=None,
) -> List[dict]:
    if extract_output_file_path is None:
        extract_output_file_path = f"presentations_text_{time.time()}.json"
        logging.info(
            "Since no extract output was specified, the output will be "
            "written to => %s", extract_output_file_path
        )
    logging.info(
        "Will start extracting presentation texts with following parameters:"
        "paths = %s, single_array_result = %b,"
        "text_fields_flattened = %b, "
        "extract_output_file_path = %s", paths, single_array_result, text_fields_flattened, extract_output_file_path
    )
    if isinstance(paths, str):
        presentations = __extract_presentation_texts_from_path__(
            paths, text_fields_flattened
        )
    else:
        presentations = [
            __extract_presentation_texts_from_path__(path, text_fields_flattened)
            for path in paths
        ]

        presentations = sum(presentations, []) if single_array_result else presentations

    with open(extract_output_file_path, "w") as output_file:
        json.dump(presentations, output_file, ensure_ascii=False)

    return presentations


def __extract_presentation_texts_from_path__(path: str, flattened: bool):
    if os.path.isdir(path):
        logging.debug(
            "%s is to a dir. Will attempt to all pptx files within it", path
        )
        presentations = __load_presentations_objects_from_dir__(path)
    else:
        presentations = __load_presentations_objects_from_file_paths__(path)
    presentations = [
        __extract_presentation_text__(presentation, flattened)
        for presentation in presentations
    ]

    return presentations


def __load_presentations_objects_from_file_paths__(file_paths):
    if isinstance(file_paths, str):
        file_paths = [file_paths]
    presentations = []
    for path in file_paths:
        if not os.path.exists(path):
            logging.error("File %s does not exist", path)
            logging.info("Will skip %s due to an error", path)
            continue
        try:
            obj = {"path": path, "presentationObj": Presentation(path)}

            presentations.append(obj)
            logging.debug("Loaded %s successfully", path)
        except ModuleNotFoundError:
            logging.error(
                "Unable to process %s . It is likely to be corrupted or incomplete. "
                "Please ensure the input is a valid pptx file", path)
            logging.info("Will skip %s due to an error", path)
    return presentations


def __load_presentations_objects_from_dir__(dir_path: str):
    paths = []
    for root, _, files in os.walk(dir_path):
        for name in files:
            if name.endswith("pptx"):
                paths.append(os.path.join(root, name))

    return __load_presentations_objects_from_file_paths__(paths)


def __extract_presentation_text__(presentation: Presentation, is_flattened: bool):
    # Collect notes
    slides_notes_texts = [
        slide.notes_slide.notes_text_frame.text
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None
        else ""
        for slide in presentation["presentationObj"].slides
    ]
    # Collect text
    slides_body_texts = [
        "".join([shape.text if shape.has_text_frame else "" for shape in slide.shapes])
        for slide in presentation["presentationObj"].slides
    ]

    result = {"path": presentation["path"]}
    if is_flattened:
        notes_text_flattened = {
            f"slide{i}NoteText": slides_notes_texts[i]
            for i in range(len(slides_notes_texts))
        }
        body_text_flattened = {
            f"slide{i}BodyText": slides_body_texts[i]
            for i in range(len(slides_body_texts))
        }
        return {
            **result,
            **body_text_flattened,
            **notes_text_flattened,
            **calculate_length_stats_for_list_of_strings(slides_body_texts, "BodyText"),
            **calculate_length_stats_for_list_of_strings(
                slides_notes_texts, "NotesText"
            ),
        }

    return {
        **result,
        "slides": [
            {"noteText": note, "bodyText": body}
            for note, body in zip(slides_notes_texts, slides_body_texts)
        ],
        "bodyTextLengthStats": calculate_length_stats_for_list_of_strings(
            slides_body_texts
        ),
        "noteTextLengthStats": calculate_length_stats_for_list_of_strings(
            slides_notes_texts
        ),
    }
